"""
DocFlow — PPTX Parser
Extracts slide text, speaker notes, tables, and images from PowerPoint files.
"""
from __future__ import annotations

from pathlib import Path
from typing import Any

from loguru import logger
from utils.markdown_utils import tables_to_markdown, clean_markdown


class PPTXParser:
    async def parse(self, file_path: Path, options: Any) -> dict:
        logger.info(f"[PPTXParser] Parsing: {file_path.name}")
        try:
            from pptx import Presentation
        except ImportError:
            return {"markdown": "", "text": "", "structured": {}, "tables": [],
                    "images": [], "metadata": {}, "hyperlinks": []}

        prs = Presentation(str(file_path))
        slides_md: list[str] = []
        tables_data: list[dict] = []
        images: list[dict] = []

        for i, slide in enumerate(prs.slides, start=1):
            slide_parts: list[str] = [f"\n## Slide {i}\n"]

            for shape in slide.shapes:
                # Text frames
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue
                        level = para.level
                        prefix = "#" * min(level + 3, 6) + " " if level > 0 else ""
                        slide_parts.append(f"{prefix}{text}")

                # Tables
                if options.extract_tables and shape.has_table:
                    rows = []
                    for row in shape.table.rows:
                        rows.append([cell.text.strip() for cell in row.cells])
                    if rows:
                        md = tables_to_markdown(rows)
                        slide_parts.append(f"\n{md}\n")
                        tables_data.append({"slide": i, "data": rows, "markdown": md})

                # Images
                if options.extract_images and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    images.append({"slide": i, "shape_name": shape.name})

            # Speaker notes
            notes = slide.notes_slide.notes_text_frame.text.strip() if slide.has_notes_slide else ""
            if notes:
                slide_parts.append(f"\n> **Notes:** {notes}\n")

            slides_md.append("\n".join(slide_parts))

        props = prs.core_properties
        metadata = {
            "title": props.title or "",
            "author": props.author or "",
            "slide_count": len(prs.slides),
        }

        full_md = "\n\n".join(slides_md)
        return {
            "markdown": clean_markdown(full_md),
            "text": "\n".join(s for s in slides_md),
            "structured": {"slide_count": len(prs.slides)},
            "tables": tables_data,
            "images": images,
            "metadata": metadata,
            "hyperlinks": [],
        }
